"""
Build the SupportBot presentation from the Elsewedy template.

The mentor's rule is that the template's design is fixed -- so this script
never restyles, moves, or recolours a single template shape. It only:
  * rewrites the text inside the template's own text boxes (reusing their
    existing runs, so font, size and colour are inherited, not re-declared), and
  * adds our own content (body copy + product screenshots) into the empty
    canvas of the white content slide, clear of every template graphic.

Structure comes from presentation/structure.pptx, which is produced by
PowerPoint itself (COM) duplicating the template slides -- that copies the
grouped freeforms and their image relationships natively, which a hand-rolled
XML deepcopy does not do reliably.

Run:
    uv run python presentation/build_deck.py
"""

from __future__ import annotations

import copy
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parent
STRUCTURE = ROOT / "structure.pptx"
# Optional CLI override, so a draft can be built and inspected while the real
# deck is open in PowerPoint (which holds an exclusive lock on it).
OUT = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "SupportBot.pptx"
SHOTS = ROOT / "screenshots"

# --- template geometry ------------------------------------------------------
# The white content slide is 28.67 x 16.00in. These bounds keep our content
# clear of: the red vertical rule (~x=1.8in), the Elsewedy mark (top right,
# below y=2.6in), the WEDY.AI mark (bottom left) and the red network graphic
# (bottom right, from roughly x=21.5in / y=11.5in).
BODY_LEFT = Inches(2.6)
BODY_WIDTH = Inches(9.0)

# The vertical band our content lives in: below the title (which can wrap to
# two lines, ending ~2.6in) and above the bottom-corner marks. Body text and
# images are CENTRED in this band rather than pinned to the top, so a slide
# reads as balanced instead of top-heavy. Text may sit a little lower than
# images because the WEDY.AI mark (bottom-left) is lower than the network
# graphic (bottom-right).
CONTENT_TOP = Inches(3.0)
CONTENT_BOTTOM_TEXT = Inches(12.6)   # left column, clears the WEDY.AI mark
IMG_BOTTOM_LIMIT = Inches(11.4)      # right column, clears the network graphic
IMG_RIGHT_LIMIT = Inches(26.6)
IMG_LEFT = Inches(12.4)

BODY_FONT = "Poppins"      # the template's own body face (see cover subtitle)
BODY_SIZE = Pt(34)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
RED = RGBColor(0xE5, 0x1B, 0x29)

# Author byline on the cover. One place to edit -- replace with the full team,
# e.g. "Built by Youssef Bassem, <name>, <name>".
BYLINE = "Built by Youssef Bassem, Mahmoud Alsayed & Somaya Ahmed"


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------
def set_lines(shape, lines: list[str]) -> None:
    """Rewrite a template text box's content, keeping its existing styling.

    Works by reusing paragraph 0 / run 0 as the style carrier and cloning it
    for extra lines, so we never restate font, size or colour -- whatever the
    template set stays set. Replacing the text frame outright would drop all
    of it and force us to re-declare the design we were told not to touch.
    """
    tf = shape.text_frame
    paras = tf.paragraphs
    template_p = copy.deepcopy(paras[0]._p)

    # Drop every paragraph after the first, and every run after the first.
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs[1:]):
        r._r.getparent().remove(r._r)

    if not first.runs:  # nothing to carry style -- nothing we can safely do
        first.text = lines[0]
    else:
        first.runs[0].text = lines[0]

    for line in lines[1:]:
        new_p = copy.deepcopy(template_p)
        # keep only the first run in the clone, then set its text
        runs = new_p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r"
        )
        for extra in runs[1:]:
            new_p.remove(extra)
        if runs:
            t = runs[0].find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
            )
            t.text = line
        tf._txBody.append(new_p)


def set_number(shape, digits: str) -> None:
    """Separator numerals: the template styles the 2nd digit red, 1st white.
    Reuse those two runs so the two-tone treatment survives."""
    runs = shape.text_frame.paragraphs[0].runs
    if len(runs) >= 2:
        runs[0].text = digits[0]
        runs[1].text = digits[1]
        for extra in runs[2:]:
            extra._r.getparent().remove(extra._r)
    else:
        set_lines(shape, [digits])


def by_name(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(f"{name!r} not found on slide (have: {[s.name for s in slide.shapes]})")


def add_body(slide, bullets: list[str], size=None) -> None:
    """Our own body copy, vertically centred in the content band.

    The text box spans the whole band and is middle-anchored, so short copy
    sits at the optical centre of the slide instead of clinging to the top.
    `size` drops a couple of points on the four-bullet slides.
    """
    size = size or BODY_SIZE
    box = slide.shapes.add_textbox(
        BODY_LEFT, CONTENT_TOP, BODY_WIDTH, Emu(CONTENT_BOTTOM_TEXT - CONTENT_TOP)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.name = BODY_FONT
        run.font.color.rgb = BODY_COLOR
        p.line_spacing = 1.35
        p.space_after = Pt(22)


def _add_shadow(pic) -> None:
    """Give a picture a soft drop shadow.

    python-pptx has no shadow API, so the effect element is added by hand. In
    a picture's spPr the sequence ends xfrm, prstGeom, [ln], effectLst -- so
    appending effectLst last is valid order. Values are EMU: ~5pt blur, ~3pt
    offset straight down (dir 5400000 = 90deg), black at 32% alpha -- enough to
    lift the screenshot off the white page without looking heavy.
    """
    spPr = pic._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    effect_lst = spPr.makeelement(qn("a:effectLst"), {})
    shadow = effect_lst.makeelement(
        qn("a:outerShdw"),
        {"blurRad": "63500", "dist": "38100", "dir": "5400000", "rotWithShape": "0"},
    )
    color = shadow.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = color.makeelement(qn("a:alpha"), {"val": "32000"})
    color.append(alpha)
    shadow.append(color)
    effect_lst.append(shadow)
    spPr.append(effect_lst)


def add_image(slide, filename: str, left=None, top=None, max_w=None, max_h=None):
    """Place a screenshot, scaled to fit and vertically centred in the image
    band, with a soft drop shadow. An explicit `top` opts out of centring
    (used when two images share a row)."""
    path = SHOTS / filename
    with Image.open(path) as im:
        ar = im.size[0] / im.size[1]

    left = IMG_LEFT if left is None else left
    max_w = (IMG_RIGHT_LIMIT - left) if max_w is None else max_w
    band_top = CONTENT_TOP if top is None else top
    avail_h = (IMG_BOTTOM_LIMIT - band_top) if max_h is None else max_h

    width = max_w
    height = Emu(int(width / ar))
    if height > avail_h:
        height = Emu(int(avail_h))
        width = Emu(int(height * ar))

    # Centre within the band unless the caller pinned a top.
    top = band_top if top is not None else Emu(int(CONTENT_TOP + (avail_h - height) / 2))

    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    _add_shadow(pic)
    return pic


def add_credits(slide, byline: str, footer: str) -> None:
    """Author credits on the dark cover, in the clear band between the template
    subtitle (ends ~10.3in) and the WEDY.AI mark (starts ~13.2in). White text,
    left-aligned to the title. Adds our own box -- no template shape is touched.
    """
    box = slide.shapes.add_textbox(Inches(2.0), Inches(10.8), Inches(14.0), Inches(1.9))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, pt) in enumerate([(byline, 30), (footer, 22)]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.name = BODY_FONT
        run.font.bold = i == 0
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.space_after = Pt(8)


def add_stats(slide, items: list[tuple[str, str]], left=Inches(15.4), top=None) -> None:
    """Big-number callouts stacked down the right half, for the data slides
    (metrics, architecture) that carry no screenshot. Each item is
    (big_number, caption): the number in brand red, the caption in grey. The
    stack is vertically centred in the image band by default, to match the
    middle-anchored body text on the left."""
    step = Inches(2.55)
    if top is None:
        stack_h = (len(items) - 1) * step + Inches(1.4)  # numbers + last caption
        top = Emu(int(CONTENT_TOP + ((IMG_BOTTOM_LIMIT - CONTENT_TOP) - stack_h) / 2))
    y = top
    for big, caption in items:
        box = slide.shapes.add_textbox(left, y, Inches(10.4), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = big
        r1.font.size = Pt(66)
        r1.font.bold = True
        r1.font.name = BODY_FONT
        r1.font.color.rgb = RED
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = caption
        r2.font.size = Pt(24)
        r2.font.name = BODY_FONT
        r2.font.color.rgb = BODY_COLOR
        y = Emu(y + step)


# ---------------------------------------------------------------------------
# Slide transitions
# ---------------------------------------------------------------------------
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def set_transition(slide, kind: str = "fade", speed: str = "med") -> None:
    """Add a slide transition.

    python-pptx has no API for this, so the element goes in by hand. Order
    inside <p:sld> is schema-enforced: cSld, clrMapOvr, transition, timing --
    inserting it anywhere else makes PowerPoint call the file corrupt.
    """
    sld = slide._element
    for existing in sld.findall(f"{{{P_NS}}}transition"):
        sld.remove(existing)

    transition = etree.SubElement(sld, f"{{{P_NS}}}transition")
    transition.set("spd", speed)
    transition.set("advClick", "1")
    etree.SubElement(transition, f"{{{P_NS}}}{kind}")

    # Re-seat it directly after clrMapOvr (or cSld) to satisfy the schema.
    anchor = sld.find(f"{{{P_NS}}}clrMapOvr")
    if anchor is None:
        anchor = sld.find(f"{{{P_NS}}}cSld")
    sld.remove(transition)
    anchor.addnext(transition)


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
def build() -> None:
    prs = Presentation(STRUCTURE)
    s = prs.slides

    # 1 -- cover -------------------------------------------------------------
    set_lines(by_name(s[0], "TextBox 17"), ["SupportBot"])
    set_lines(
        by_name(s[0], "TextBox 18"),
        ["An AI support assistant for", "Elsewedy field engineers"],
    )
    add_credits(
        s[0],
        BYLINE,
        "AI Section — Elsewedy Electric | Summer Internship 2026",
    )

    # 2 -- the problem -------------------------------------------------------
    # Numbers here are real, measured from data/cases_clean.jsonl: 22 of 66
    # searchable cases repeat a problem already solved on another unit (33%),
    # and one problem recurs across three products. The time-cost and
    # knowledge-loss points are framed qualitatively rather than with an
    # invented figure.
    set_lines(by_name(s[1], "TextBox 6"), ["The problem"])
    add_body(
        s[1],
        [
            "A device fails in the field. The fix almost always exists "
            "already — in an old ticket, a senior engineer's memory, a "
            "colleague's inbox. Reaching it means asking around and waiting.",
            "These faults repeat. In our own case base, one resolved case in "
            "three is a problem that had already been solved on another unit — "
            "one recurs across three different products.",
            "And the knowledge lives in people. When an experienced engineer "
            "leaves, the fixes they carried leave with them.",
            "SupportBot turns those resolved cases into a shared memory any "
            "engineer can search in seconds.",
        ],
        size=Pt(30),
    )
    add_stats(
        s[1],
        [("1 in 3", "resolved cases repeat a problem\nalready solved elsewhere")],
    )

    # 3 -- separator 01 ------------------------------------------------------
    set_number(by_name(s[2], "TextBox 8"), "01")
    set_lines(by_name(s[2], "TextBox 7"), ["The Product"])

    # 4 -- what it does -----------------------------------------------------
    set_lines(by_name(s[3], "TextBox 6"), ["Find the fix"])
    add_body(
        s[3],
        [
            "Engineers describe a problem the way they'd say it out loud — no "
            "error codes, no lookup tables.",
            "SupportBot searches 66 resolved cases spanning 7 products and 8 "
            "issue categories, by meaning rather than by keyword.",
            "Type only a product name and it reports how many past cases exist "
            "for it, with an example to start from.",
            "Suggested questions get you started in a single click.",
        ],
        size=Pt(30),
    )
    add_image(s[3], "slide-03-landing-hero-light.png")

    # 5 -- citation ----------------------------------------------------------
    set_lines(by_name(s[4], "TextBox 6"), ["Cited answers"])
    add_body(
        s[4],
        [
            "Every answer names the exact past case behind it — Case ID, the "
            "original problem, and its resolution.",
            "The fix shown is the fix that was actually recorded, not a "
            "paraphrase of it.",
            "Answers can be copied, or rated helpful or not helpful to flag "
            "where the knowledge base needs a better case.",
        ],
    )
    add_image(s[4], "slide-09-grounded-answer.png")

    # 6 -- conversation memory -----------------------------------------------
    set_lines(by_name(s[5], "TextBox 6"), ["It remembers"])
    add_body(
        s[5],
        [
            "Follow-up questions work: “does it happen on the FlowMeter X100 "
            "too” is understood in the context of what came before.",
            "The thread is kept, so an engineer never has to restate the "
            "problem.",
            "Every reply is still grounded in its own cited case.",
        ],
    )
    add_image(s[5], "slide-10-followup-memory.png")

    # 7 -- explainability ----------------------------------------------------
    set_lines(by_name(s[6], "TextBox 6"), ["Shows its work"])
    add_body(
        s[6],
        [
            "“See how this was found” reveals the search space behind the "
            "answer.",
            "Each point is a past case; the highlighted ones are what the "
            "question actually matched.",
            "Engineers can see why an answer was chosen — not just be asked to "
            "trust it.",
        ],
    )
    add_image(s[6], "slide-11-embedding-map.png")

    # 8 -- separator 02 ------------------------------------------------------
    set_number(by_name(s[7], "TextBox 8"), "02")
    set_lines(by_name(s[7], "TextBox 7"), ["Trust"])

    # 9 -- scope -------------------------------------------------------------
    set_lines(by_name(s[8], "TextBox 6"), ["Stays in scope"])
    add_body(
        s[8],
        [
            "General-knowledge questions are declined — this is a product "
            "support tool, not a general chatbot.",
            "If nothing in the knowledge base matches, it says so instead of "
            "guessing an answer.",
            "An unrelated case is never dressed up as a resolution.",
        ],
    )
    add_image(s[8], "slide-08-threshold-rejected.png")

    # 10 -- small talk -------------------------------------------------------
    set_lines(by_name(s[9], "TextBox 6"), ["Small talk"])
    add_body(
        s[9],
        [
            "Greetings and casual questions get a warm, natural reply.",
            "It answers, then steers back to what it can actually help with.",
            "Small talk never produces a fabricated Case ID — the citation "
            "format appears only when there is a real case behind it.",
        ],
    )
    add_image(s[9], "slide-12-small-talk.png")

    # 11 -- graceful degradation ---------------------------------------------
    set_lines(by_name(s[10], "TextBox 6"), ["Always answers"])
    add_body(
        s[10],
        [
            "If the answer-writing service is unavailable, the matching cases "
            "are still listed in full.",
            "The engineer still gets Case ID, problem and resolution — only "
            "the written summary is missing, and it says so.",
            "It degrades gracefully rather than showing an error screen.",
        ],
    )
    add_image(s[10], "slide-13-offline-fallback.png")

    # 12 -- how we measured it (metrics) -------------------------------------
    # Every number here is reproduced by `uv run python -m eval.eval_retriever`.
    set_lines(by_name(s[11], "TextBox 6"), ["How we measured it"])
    add_body(
        s[11],
        [
            "A 15-query Gold Set — 5 exact-identifier, 7 paraphrase, 3 "
            "out-of-domain — each hand-labelled with its correct case.",
            "Hit@1, Hit@3, Hit@5 and MRR@5 against those labels — re-run in CI "
            "on every commit, so the score is reproducible, not a one-off.",
            "The abstention cutoff is derived, not guessed: answerable cases "
            "land at distance ≤ 0.582, out-of-domain at ≥ 0.874. We set it at "
            "0.735 — inside that clean, non-overlapping gap, so every "
            "answerable case falls below it and every out-of-domain one above.",
        ],
        size=Pt(30),
    )
    add_stats(
        s[11],
        [
            ("100%", "Hit@1 on every answerable query"),
            ("1.000", "MRR@5 — the right case ranked first"),
            ("100%", "correct abstention on out-of-domain"),
        ],
    )

    # 13 -- why hybrid search (architecture trade-off) -----------------------
    set_lines(by_name(s[12], "TextBox 6"), ["Why hybrid search"])
    add_body(
        s[12],
        [
            "Meaning-search alone already scores 100% here — it handles "
            "paraphrase well. But it blurs exact model numbers: G2 and G3 sit "
            "almost on top of each other in vector space.",
            "Keyword-search keeps those identifiers distinct, yet on its own "
            "misses paraphrase — 71% Hit@1.",
            "We run both and fuse by rank (RRF). Hybrid keeps each one's "
            "strength and trades away neither — we run it for exact-identifier "
            "robustness, not because fusion beats meaning-search on this set.",
        ],
        size=Pt(30),
    )
    add_stats(
        s[12],
        [
            ("100%", "meaning-only Hit@1 — strong on paraphrase"),
            ("71%", "keyword-only Hit@1 — misses paraphrase"),
            ("100%", "hybrid — keeps both strengths"),
        ],
    )

    # 14 -- separator 03 -----------------------------------------------------
    set_number(by_name(s[13], "TextBox 8"), "03")
    set_lines(by_name(s[13], "TextBox 7"), ["The Experience"])

    # 15 -- bilingual --------------------------------------------------------
    set_lines(by_name(s[14], "TextBox 6"), ["Fully bilingual"])
    add_body(
        s[14],
        [
            "One toggle switches the entire interface between English and "
            "Arabic — including a true right-to-left layout — and the choice "
            "is remembered next time.",
            "Ask a question in Arabic and the answer comes back in Arabic.",
            "Arabic is set in a dedicated typeface, not left to a browser "
            "fallback.",
        ],
    )
    add_image(s[14], "slide-15-arabic-rtl.png")

    # 16 -- theming + accessibility ------------------------------------------
    set_lines(by_name(s[15], "TextBox 6"), ["Light and dark"])
    add_body(
        s[15],
        [
            "One click re-themes the whole application; it follows the "
            "operating system by default and remembers an explicit choice.",
            "Every text and surface pairing was contrast-checked for "
            "accessibility in both themes.",
            "Keyboard focus is visible throughout, and animation is reduced "
            "for anyone whose system asks for that.",
        ],
    )
    add_image(s[15], "slide-16-dark-mode.png")

    # 17 -- responsive / mobile (its own slide, with a real mobile chat shot) -
    set_lines(by_name(s[16], "TextBox 6"), ["On any screen"])
    add_body(
        s[16],
        [
            "The full experience runs down to phone width — where field "
            "support actually happens.",
            "Chat, cited answers and the case view reflow to a single column; "
            "nothing is cut off.",
            "The sidebar collapses into a slide-in drawer, reached from the "
            "menu button top-left, that opens from the correct side in Arabic.",
        ],
    )
    # A single real mobile-chat screenshot (captured live on a phone-width
    # viewport, demo account), centred in the right half with a drop shadow.
    add_image(s[16], "slide-19-mobile-chat.png", left=Inches(17.5), max_w=Inches(3.9))

    # 18 -- accounts ---------------------------------------------------------
    set_lines(by_name(s[17], "TextBox 6"), ["Your account"])
    add_body(
        s[17],
        [
            "Engineers create an account and stay signed in — a refresh never "
            "signs you out mid-job.",
            "A forgotten password is reset with a code sent to your email.",
            "Sign-in is protected against repeated password guessing.",
        ],
    )
    add_image(s[17], "slide-14-auth-signin.png")

    # 19 -- thank you --------------------------------------------------------
    set_lines(by_name(s[18], "TextBox 17"), ["Thank you"])

    # transitions ------------------------------------------------------------
    # Fade throughout, with a push on the three section separators so a section
    # change reads as a change. Deliberately restrained: this is a template
    # deck, not a showreel.
    for i, slide in enumerate(s, start=1):
        set_transition(slide, "push" if i in (3, 8, 14) else "fade")

    prs.save(OUT)
    print(f"[deck] wrote {OUT}  ({len(s.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
